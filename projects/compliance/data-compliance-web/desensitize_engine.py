from __future__ import annotations

import csv
import importlib.util
import json
import os
import re
import shutil
import subprocess
import sys
from collections import Counter, defaultdict
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

from docx import Document

try:
    from scripts.ocr_text import OcrUnavailable, ensure_ocr_available, ocr_image_to_data
except ModuleNotFoundError:  # Allows importing as web.desensitize_engine in checks.
    from web.scripts.ocr_text import OcrUnavailable, ensure_ocr_available, ocr_image_to_data

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None

try:
    import fitz  # PyMuPDF
except Exception:  # pragma: no cover - optional dependency
    fitz = None

try:
    import pandas as pd
except Exception:  # pragma: no cover - optional dependency
    pd = None

try:
    from PIL import Image, ImageDraw
except Exception:  # pragma: no cover - optional dependency
    Image = None
    ImageDraw = None

try:
    from presidio_analyzer import AnalyzerEngine
except Exception:  # pragma: no cover - optional dependency
    AnalyzerEngine = None


def _ensure_redaction_in_path() -> None:
    """将 legalwork 根目录加入 sys.path，以便导入 redaction 包。"""
    current = Path(__file__).resolve()
    for parent in current.parents:
        if (parent / 'redaction' / '__init__.py').exists():
            parent_str = str(parent)
            if parent_str not in sys.path:
                sys.path.insert(0, parent_str)
            break


_ensure_redaction_in_path()

try:
    from redaction.detector import RedactionDetector

    REDACTION_AVAILABLE = True
except Exception as _redaction_import_error:  # pragma: no cover - optional dependency
    RedactionDetector = None  # type: ignore[misc, assignment]
    REDACTION_AVAILABLE = False

LEGAL_SUBJECT_TYPES = ['person_name', 'company_name']
FAST_SUBJECT_TYPES = ['company_name']
REDACTION_FILL = (229, 231, 235)
PDF_REDACTION_FILL = tuple(channel / 255 for channel in REDACTION_FILL)
REDACTION_PADDING = 1.5

PERSON_CONTEXT_PATTERN = re.compile(
    r'(?:原告|被告|第三人|上诉人|被上诉人|申请人|被申请人|申请执行人|被执行人|'
    r'甲方|乙方|丙方|丁方|委托人|受托人|法定代表人|负责人|联系人|姓名)'
    r'[：:\s，,、]*([\u4e00-\u9fa5]{2,4})'
)
ENGLISH_PERSON_CONTEXT_PATTERN = re.compile(
    r'\b(?:Legal\s+Representative|Representative|Contact(?:s)?|Director|Manager|Name)'
    r'[ \t]*[:：]?[ \t]*([A-Z][A-Za-z.\'-]+(?:[ \t]+[A-Z][A-Za-z.\'-]+){1,3})\b'
)
ENGLISH_LABEL_CHINESE_PERSON_PATTERN = re.compile(
    r'\b(?:Legal\s+Representative|Representative|Contact(?:s)?|Director|Manager|Name)'
    r'[ \t]*[:：]?[ \t]*([\u4e00-\u9fa5]{2,4})'
)
PERSON_LIST_PATTERN = re.compile(
    r'(?<![\u4e00-\u9fa5])([\u4e00-\u9fa5]{2,4})(?:、|和|与)([\u4e00-\u9fa5]{2,4})(?![\u4e00-\u9fa5])'
)
ENGLISH_COMPANY_PATTERN = re.compile(
    r'\b([A-Z][A-Za-z0-9&.\',()/-]*(?:\s+[A-Z][A-Za-z0-9&.\',()/-]*){0,12}\s+'
    r'(?:Co\.?,?\s*Ltd\.?|Company|Corporation|Corp\.?|Limited|LLC|Inc\.?|Holdings?|Group|'
    r'Finance\s+Co\.?,?\s*Ltd\.?))\b',
    re.IGNORECASE,
)
BRANDING_TEXT_PATTERN = re.compile(
    r'\b(?:ORIGINATOR|SERVICER|LOGO)\b|\[[^\]]*(?:ORIGINATOR|LOGO)[^\]]*\]',
    re.IGNORECASE,
)
ENGLISH_LEGAL_NORM_FOLLOWING_PATTERN = re.compile(r'^\s+(?:Law|Act|Code|Regulation|Regulations|Rules)\b', re.IGNORECASE)


@dataclass(frozen=True)
class SubjectMapping:
    entity_type: str
    original: str
    redacted: str
    location: str
    confidence: float

    def to_dict(self) -> dict[str, Any]:
        return {
            'entity_type': self.entity_type,
            'original': self.original,
            'redacted': self.redacted,
            'location': self.location,
            'confidence': round(self.confidence, 3),
        }


TEXT_EXTENSIONS = {
    '.txt',
    '.md',
    '.markdown',
    '.log',
    '.rtf',
    '.html',
    '.htm',
    '.xml',
    '.yaml',
    '.yml',
    '.toml',
    '.ini',
    '.cfg',
    '.conf',
    '.env',
}
DOC_EXTENSIONS = {'.docx', '.doc'}
PDF_EXTENSIONS = {'.pdf'}
TABLE_EXTENSIONS = {'.csv', '.tsv', '.xlsx', '.xls', '.ods'}
JSON_EXTENSIONS = {'.json', '.jsonl', '.ndjson'}
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp', '.bmp', '.tif', '.tiff'}
PRESENTATION_EXTENSIONS = {'.pptx'}


@dataclass(frozen=True)
class Finding:
    entity_type: str
    start: int
    end: int
    score: float
    replacement: str
    surface: str
    locator: str
    preview: str

    def to_dict(self) -> dict[str, Any]:
        return {
            'entity_type': self.entity_type,
            'start': self.start,
            'end': self.end,
            'score': round(self.score, 3),
            'replacement': self.replacement,
            'surface': self.surface,
            'locator': self.locator,
            'preview': self.preview,
        }


@dataclass(frozen=True)
class ImageRedactionResult:
    output_file: Path
    findings: list[Finding]
    subjects: list[SubjectMapping]
    preview_text: str


def _looks_like_person_name(value: str) -> bool:
    if not 2 <= len(value) <= 4:
        return False
    noise = {'公司', '企业', '集团', '法院', '银行', '合同', '协议', '数据', '用户', '个人', '信息'}
    return value not in noise and not any(word in value for word in noise)


def _detect_person_subjects(text: str) -> list[Any]:
    """用轻量规则识别常见中文自然人姓名，避免为人名加载重型 NER。"""
    entities: list[Any] = []
    if not text:
        return entities

    class _Entity:
        def __init__(self, original: str, start: int, end: int, confidence: float) -> None:
            self.entity_type = 'person_name'
            self.text = original
            self.start = start
            self.end = end
            self.confidence = confidence

    for match in PERSON_CONTEXT_PATTERN.finditer(text):
        original = match.group(1)
        if _looks_like_person_name(original):
            entities.append(_Entity(original, match.start(1), match.end(1), 0.82))

    for match in ENGLISH_PERSON_CONTEXT_PATTERN.finditer(text):
        original = match.group(1).strip()
        if original and not ENGLISH_COMPANY_PATTERN.search(original):
            entities.append(_Entity(original, match.start(1), match.end(1), 0.74))

    for match in ENGLISH_LABEL_CHINESE_PERSON_PATTERN.finditer(text):
        original = match.group(1).strip()
        if _looks_like_person_name(original):
            entities.append(_Entity(original, match.start(1), match.end(1), 0.8))

    for match in PERSON_LIST_PATTERN.finditer(text):
        for group_index in (1, 2):
            original = match.group(group_index)
            if _looks_like_person_name(original):
                entities.append(_Entity(original, match.start(group_index), match.end(group_index), 0.72))

    return entities


def _detect_english_company_subjects(text: str) -> list[Any]:
    """识别英文公司名，覆盖封面/表格中常见的 Co., Ltd. 等主体。"""
    entities: list[Any] = []
    if not text:
        return entities

    class _Entity:
        def __init__(self, original: str, start: int, end: int, confidence: float) -> None:
            self.entity_type = 'company_name'
            self.text = re.sub(r'\s+', ' ', original).strip(' ,;')
            self.start = start
            self.end = end
            self.confidence = confidence

    for match in ENGLISH_COMPANY_PATTERN.finditer(text):
        original = match.group(1)
        if len(original.strip()) >= 5:
            entities.append(_Entity(original, match.start(1), match.end(1), 0.78))
    return entities


def _is_public_legal_norm_subject(text: str, entity: Any) -> bool:
    """公开法律规范名称不作为案件主体替换。"""
    start = int(getattr(entity, 'start', 0))
    end = int(getattr(entity, 'end', 0))
    original = str(getattr(entity, 'text', ''))
    if not original or end <= start:
        return False

    following = text[end:end + 24]
    if ENGLISH_LEGAL_NORM_FOLLOWING_PATTERN.match(following):
        return True

    left = text.rfind('《', 0, start + 1)
    right = text.find('》', end)
    if left >= 0 and right >= 0 and text.find('》', left, start) < 0:
        title = text[left + 1:right]
        if re.search(r'(法|法律|法典|条例|规定|办法|规则|细则|解释|决定|意见|通知)$', title):
            return True
    return False


def detect_legal_subject_entities(
    text: str,
    *,
    detector: Any | None = None,
) -> list[Any]:
    """返回带位置的法律主体实体，供文本替换和版式遮盖共用。"""
    entities: list[Any] = []
    if REDACTION_AVAILABLE and detector is not None:
        entities.extend(detector.detect(text, entity_types=FAST_SUBJECT_TYPES, use_semantic=False))
    entities.extend(_detect_english_company_subjects(text))
    entities.extend(_detect_person_subjects(text))
    entities = [entity for entity in entities if not _is_public_legal_norm_subject(text, entity)]
    return dedupe_subject_entities(entities)


def dedupe_subject_entities(entities: list[Any]) -> list[Any]:
    ordered = sorted(
        entities,
        key=lambda item: (
            int(getattr(item, 'start', 0)),
            -(int(getattr(item, 'end', 0)) - int(getattr(item, 'start', 0))),
            -float(getattr(item, 'confidence', 0.0)),
        ),
    )
    kept: list[Any] = []
    for entity in ordered:
        start = int(getattr(entity, 'start', 0))
        end = int(getattr(entity, 'end', 0))
        if end <= start:
            continue
        if any(start < int(getattr(current, 'end', 0)) and end > int(getattr(current, 'start', 0)) for current in kept):
            continue
        kept.append(entity)
    return sorted(kept, key=lambda item: int(getattr(item, 'start', 0)))


def redact_legal_subjects(
    text: str,
    locator: str = '',
    *,
    detector: Any | None = None,
) -> tuple[str, list[SubjectMapping]]:
    """识别并替换法律主体（自然人、公司），返回替换后文本与可逆映射表。"""
    if not text or not text.strip():
        return text, []

    entities = detect_legal_subject_entities(text, detector=detector)
    if not entities:
        return text, []

    # 按实体文本聚类，保证同一主体对应同一 token
    clusters: dict[str, dict[str, Any]] = {}
    cluster_counters: dict[str, int] = {'person_name': 0, 'company_name': 0}
    token_labels = {
        'person_name': ['当事人甲', '当事人乙', '当事人丙', '当事人丁', '当事人戊', '当事人己', '当事人庚', '当事人辛', '当事人壬', '当事人癸'],
        'company_name': ['A公司', 'B公司', 'C公司', 'D公司', 'E公司', 'F公司', 'G公司', 'H公司', 'I公司', 'J公司'],
    }

    for entity in entities:
        entity_type = getattr(entity, 'entity_type', '')
        original = getattr(entity, 'text', '')
        if entity_type not in token_labels or not original:
            continue
        if original not in clusters:
            cluster_counters[entity_type] += 1
            label_index = cluster_counters[entity_type] - 1
            if label_index < len(token_labels[entity_type]):
                token = token_labels[entity_type][label_index]
            else:
                token = f'{token_labels[entity_type][-1].replace("公司", "").replace("当事人", "")}_{cluster_counters[entity_type]}'
                if entity_type == 'company_name':
                    token = f'{token}公司'
                else:
                    token = f'当事人{token}'
            clusters[original] = {
                'entity_type': entity_type,
                'token': token,
            }

    if not clusters:
        return text, []

    # 从后向前替换，避免位置偏移
    sorted_entities = sorted(
        entities,
        key=lambda e: (getattr(e, 'start', 0), getattr(e, 'end', 0)),
        reverse=True,
    )
    redacted_text = text
    subject_mappings: list[SubjectMapping] = []
    seen: set[tuple[int, int]] = set()

    for entity in sorted_entities:
        start = getattr(entity, 'start', 0)
        end = getattr(entity, 'end', 0)
        original = getattr(entity, 'text', '')
        entity_type = getattr(entity, 'entity_type', '')
        if start < 0 or end > len(text) or end <= start:
            continue
        if (start, end) in seen:
            continue
        # 跳过与已处理区域重叠的项
        if any((start < existing_end and end > existing_start) for existing_start, existing_end in seen):
            continue
        seen.add((start, end))
        if original not in clusters:
            continue
        token = clusters[original]['token']
        redacted_text = redacted_text[:start] + token + redacted_text[end:]
        subject_mappings.append(SubjectMapping(
            entity_type=entity_type,
            original=original,
            redacted=token,
            location=locator,
            confidence=float(getattr(entity, 'confidence', 0.0)),
        ))

    return redacted_text, subject_mappings


def sanitize_text_and_subjects(
    text: str,
    engine: Desensitizer,
    *,
    surface: str = 'text',
    locator: str = '',
) -> tuple[str, list[Finding], list[SubjectMapping]]:
    """先执行隐私信息脱敏，再对法律主体进行可逆替换。"""
    sanitized, findings = engine.sanitize_text(text, surface=surface, locator=locator)
    redacted, subject_mappings = redact_legal_subjects(
        sanitized,
        locator=locator,
        detector=engine._subject_detector,
    )
    return redacted, findings, subject_mappings


REGEX_RULES: list[tuple[str, re.Pattern[str], float]] = [
    (
        'EMAIL_ADDRESS',
        re.compile(r'(?<![\w.+-])[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}(?![\w.-])'),
        0.95,
    ),
    (
        'PHONE_NUMBER',
        re.compile(r'(?<!\d)(?:\+?86[-\s]?)?1[3-9]\d[-\s]?\d{4}[-\s]?\d{4}(?!\d)'),
        0.96,
    ),
    (
        'ID_CARD',
        re.compile(r'(?<![0-9A-Za-z])\d{6}(?:18|19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dXx](?![0-9A-Za-z])'),
        0.96,
    ),
    (
        'BANK_CARD',
        re.compile(r'(?<!\d)\d(?:[ -]?\d){14,18}(?!\d)'),
        0.8,
    ),
    (
        'IP_ADDRESS',
        re.compile(r'(?<!\d)(?:25[0-5]|2[0-4]\d|1?\d?\d)(?:\.(?:25[0-5]|2[0-4]\d|1?\d?\d)){3}(?!\d)'),
        0.9,
    ),
    (
        'API_KEY',
        re.compile(r'(?i)\b(?:api[_-]?key|secret|token|password|passwd|access[_-]?key)\s*[:=]\s*[\'"][^\'"]{8,}[\'"]'),
        0.9,
    ),
    (
        'PRIVATE_KEY',
        re.compile(r'-----BEGIN (?:RSA |EC |OPENSSH |)?PRIVATE KEY-----[\s\S]+?-----END (?:RSA |EC |OPENSSH |)?PRIVATE KEY-----'),
        0.99,
    ),
    (
        'ADDRESS',
        re.compile(r'(?:[\u4e00-\u9fa5]{2,}(?:省|自治区|市|区|县|镇|乡|街道)[\u4e00-\u9fa5A-Za-z0-9\-]{0,24})|(?:[\u4e00-\u9fa5]{2,}(?:路|街|巷|弄)\d*号?[\u4e00-\u9fa5A-Za-z0-9\-]{0,12})'),
        0.72,
    ),
]

PRESIDIO_ENTITY_MAP = {
    'EMAIL_ADDRESS': 'EMAIL_ADDRESS',
    'PHONE_NUMBER': 'PHONE_NUMBER',
    'CREDIT_CARD': 'BANK_CARD',
    'IP_ADDRESS': 'IP_ADDRESS',
    'PERSON': 'PERSON',
    'LOCATION': 'ADDRESS',
    'US_SSN': 'ID_NUMBER',
    'IBAN_CODE': 'BANK_CARD',
    'CRYPTO': 'API_KEY',
}


class Desensitizer:
    def __init__(self) -> None:
        self._analyzer: Any | None = None
        self._subject_detector: Any | None = None
        self.presidio_available = False
        self.subject_detector_available = False
        use_presidio = os.environ.get('COMPLIANCEAI_USE_PRESIDIO', '').lower() in {'1', 'true', 'yes'}
        if use_presidio and AnalyzerEngine is not None and importlib.util.find_spec('en_core_web_sm') is not None:
            try:
                self._analyzer = AnalyzerEngine()
                self.presidio_available = True
            except Exception:
                self._analyzer = None
        if REDACTION_AVAILABLE:
            try:
                self._subject_detector = RedactionDetector()
                self.subject_detector_available = True
            except Exception:
                self._subject_detector = None

    def sanitize_text(self, text: str, *, surface: str = 'text', locator: str = '') -> tuple[str, list[Finding]]:
        findings = self._detect(text, surface=surface, locator=locator)
        sanitized = replace_spans(text, findings)
        return sanitized, findings

    def _detect(self, text: str, *, surface: str, locator: str) -> list[Finding]:
        findings: list[Finding] = []
        for entity_type, pattern, score in REGEX_RULES:
            for match in pattern.finditer(text):
                value = match.group(0)
                if entity_type == 'BANK_CARD' and len(re.sub(r'\D', '', value)) < 15:
                    continue
                if entity_type == 'BANK_CARD' and looks_like_id_card(value):
                    continue
                findings.append(make_finding(entity_type, match.start(), match.end(), score, surface, locator, value))

        if self._analyzer is not None and text.strip():
            try:
                for result in self._analyzer.analyze(text=text[:100_000], language='en'):
                    entity_type = PRESIDIO_ENTITY_MAP.get(result.entity_type, result.entity_type)
                    value = text[result.start:result.end]
                    findings.append(
                        make_finding(
                            entity_type,
                            result.start,
                            result.end,
                            float(result.score),
                            surface,
                            locator,
                            value,
                        )
                    )
            except Exception:
                pass

        return dedupe_findings(findings)


def make_finding(
    entity_type: str,
    start: int,
    end: int,
    score: float,
    surface: str,
    locator: str,
    value: str,
) -> Finding:
    return Finding(
        entity_type=entity_type,
        start=start,
        end=end,
        score=score,
        replacement=desensitize_value(entity_type, value),
        surface=surface,
        locator=locator,
        preview=mask_preview(value),
    )


def dedupe_findings(findings: list[Finding]) -> list[Finding]:
    ordered = sorted(findings, key=lambda item: (item.start, -(item.end - item.start), -item.score))
    kept: list[Finding] = []
    for finding in ordered:
        if any(not (finding.end <= current.start or finding.start >= current.end) for current in kept):
            continue
        kept.append(finding)
    return sorted(kept, key=lambda item: item.start)


def replace_spans(text: str, findings: list[Finding]) -> str:
    output: list[str] = []
    cursor = 0
    for finding in sorted(findings, key=lambda item: item.start):
        if finding.start < cursor:
            continue
        output.append(text[cursor:finding.start])
        output.append(finding.replacement)
        cursor = finding.end
    output.append(text[cursor:])
    return ''.join(output)


def desensitize_value(entity_type: str, value: str) -> str:
    compact = re.sub(r'\s+', '', value)
    if entity_type == 'PHONE_NUMBER':
        digits = re.sub(r'\D', '', value)
        if digits.startswith('86') and len(digits) == 13:
            digits = digits[2:]
        if len(digits) >= 7:
            return f'{digits[:3]}****{digits[-4:]}'
        return mask_middle(compact, left=2, right=2)

    if entity_type == 'EMAIL_ADDRESS' and '@' in value:
        local, domain = value.split('@', 1)
        visible = local[:2] if len(local) > 2 else local[:1]
        return f'{visible}***@{domain}'

    if entity_type in {'ID_CARD', 'ID_NUMBER'}:
        return mask_middle(compact, left=3, right=4)

    if entity_type == 'BANK_CARD':
        digits = re.sub(r'\D', '', value)
        return mask_middle(digits, left=6, right=4)

    if entity_type == 'IP_ADDRESS':
        parts = value.split('.')
        if len(parts) == 4:
            return f'{parts[0]}.***.***.{parts[-1]}'
        return mask_middle(compact, left=2, right=2)

    if entity_type == 'ADDRESS':
        return mask_middle(value.strip(), left=3, right=3)

    if entity_type == 'PERSON':
        return value[:1] + '*' * max(1, len(value) - 1)

    if entity_type in {'API_KEY', 'PRIVATE_KEY'}:
        return f'[{entity_type}]'

    return mask_middle(compact, left=2, right=2)


def mask_middle(value: str, *, left: int, right: int) -> str:
    if not value:
        return ''
    if len(value) <= left + right:
        if len(value) <= 2:
            return value[:1] + '*'
        return value[:1] + '*' * (len(value) - 2) + value[-1:]
    return value[:left] + '*' * (len(value) - left - right) + value[-right:]


def mask_preview(value: str) -> str:
    compact = re.sub(r'\s+', ' ', value).strip()
    if not compact:
        return ''
    if len(compact) <= 4:
        return '*' * len(compact)
    return f'{compact[:2]}***{compact[-2:]}'


def looks_like_id_card(value: str) -> bool:
    compact = re.sub(r'\s+', '', value)
    return bool(re.fullmatch(r'\d{6}(?:18|19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dXx]', compact))


def read_text_file(path: Path) -> str:
    if path.suffix.lower() in {'.rtf', '.doc'} and shutil.which('textutil'):
        run = subprocess.run(
            ['textutil', '-convert', 'txt', '-stdout', str(path)],
            capture_output=True,
            text=True,
            check=False,
        )
        if run.returncode == 0 and run.stdout.strip():
            return run.stdout

    data = path.read_bytes()
    for encoding in ('utf-8', 'utf-8-sig', 'gb18030', 'latin-1'):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode('utf-8', errors='replace')


def text_output_suffix(input_path: Path) -> str:
    suffix = input_path.suffix.lower()
    if suffix in {'.rtf', '.doc'}:
        return '.txt'
    return suffix or '.txt'


def build_report(
    *,
    task_id: str,
    document_name: str,
    input_name: str,
    input_type: str,
    output_file: Path | None,
    findings: list[Finding],
    warnings: list[str],
    engine: Desensitizer,
) -> dict[str, Any]:
    counts = Counter(item.entity_type for item in findings)
    surfaces = defaultdict(int)
    for item in findings:
        surfaces[item.surface] += 1

    return {
        'task_id': task_id,
        'document_name': document_name,
        'input_name': input_name,
        'input_type': input_type,
        'status': 'completed',
        'strategy': 'format_preserving_mask',
        'engine': {
            'presidio_available': engine.presidio_available,
            'presidio_default_enabled': os.environ.get('COMPLIANCEAI_USE_PRESIDIO', '').lower() in {'1', 'true', 'yes'},
            'custom_chinese_rules_enabled': True,
            'legal_subject_detector_available': engine.subject_detector_available,
        },
        'summary': {
            'total_findings': len(findings),
            'entity_counts': dict(sorted(counts.items())),
            'surface_counts': dict(sorted(surfaces.items())),
        },
        'output': {
            'file_name': output_file.name if output_file else '',
            'relative_name': output_file.name if output_file else '',
        },
        'findings': [item.to_dict() for item in findings[:500]],
        'warnings': warnings,
        'residual_risk': '自动脱敏不能保证识别全部敏感信息，正式外发前仍建议抽样复核。',
    }


def render_report_markdown(report: dict[str, Any]) -> str:
    lines = [
        f'# {report.get("document_name", "数据脱敏处理")} 脱敏报告',
        '',
        f'- 输入文件：{report.get("input_name", "")}',
        f'- 输入类型：{report.get("input_type", "")}',
        '- 脱敏策略：保留格式打码',
        f'- 命中总数：{report.get("summary", {}).get("total_findings", 0)}',
        '',
        '## 命中类型统计',
    ]
    counts = report.get('summary', {}).get('entity_counts', {})
    if counts:
        for key, value in counts.items():
            lines.append(f'- {key}: {value}')
    else:
        lines.append('- 未命中敏感信息')
    lines.extend(['', '## 处理说明'])
    lines.append(report.get('residual_risk', ''))
    warnings = report.get('warnings') or []
    if warnings:
        lines.extend(['', '## 注意事项'])
        for warning in warnings:
            lines.append(f'- {warning}')
    return '\n'.join(lines).rstrip() + '\n'


def write_retention_note(path: Path) -> None:
    path.write_text(
        '\n'.join(
            [
                '原始文件仅保存在本地任务目录，用于完成本次数据脱敏处理。',
                '系统默认输出保留格式打码后的脱敏文件和处理报告。',
                '自动识别存在漏检和误检风险，正式外发前请对脱敏结果进行抽样复核。',
                '',
            ]
        ),
        encoding='utf-8',
    )


def render_subject_mapping_markdown(
    task_id: str,
    document_name: str,
    mappings: list[SubjectMapping],
) -> str:
    """生成主体逆向映射 Markdown 文档。"""
    lines = [
        '# 主体逆向映射表',
        '',
        f'- 任务编号：{task_id}',
        f'- 材料名称：{document_name or "未命名材料"}',
        f'- 生成时间：{datetime.now().isoformat()}',
        f'- 主体数量：{len(mappings)}',
        '',
        '## 说明',
        '',
        '本文件记录材料脱敏中影响案件法律主要事实的主体（自然人、公司等）的原始值与脱敏值的对应关系。',
        '手机号、身份证号、银行卡号、地址、邮箱等个人隐私信息已完全脱敏，不在本表中。',
        '',
        '## 主体映射',
        '',
        '| 原始值 | 脱敏值 | 实体类型 | 位置 | 置信度 |',
        '|--------|--------|----------|------|--------|',
    ]
    for m in mappings:
        entity_type_label = '人名' if m.entity_type == 'person_name' else '公司名称'
        lines.append(
            f'| {m.original} | {m.redacted} | {entity_type_label} | {m.location} | {m.confidence} |'
        )
    lines.extend([
        '',
        '## 还原说明',
        '',
        '如需还原原始文本，请将上表中的「脱敏值」替换回「原始值」。同一主体在原文中可能出现多次，均已统一替换。',
        '',
    ])
    return '\n'.join(lines)


def build_subject_mapping_json(
    task_id: str,
    document_name: str,
    mappings: list[SubjectMapping],
) -> dict[str, Any]:
    """生成主体逆向映射 JSON 数据。"""
    return {
        'task_id': task_id,
        'document_name': document_name or '未命名材料',
        'generated_at': datetime.now().isoformat(),
        'subject_count': len(mappings),
        'legal_subjects': [m.to_dict() for m in mappings],
        'notes': (
            '本文件记录影响案件法律主要事实的主体（自然人、公司等）的原始值与脱敏值对应关系。'
            '手机号、身份证号、银行卡号、地址、邮箱等个人隐私信息已完全脱敏，不在本表中。'
        ),
    }


def copy_outputs_to_directory(
    output_dir: Path,
    output_file: Path,
    subject_mapping_md: Path,
    subject_mapping_json: Path,
    document_name: str,
) -> None:
    """将脱敏产物复制到用户指定的输出目录。"""
    output_dir.mkdir(parents=True, exist_ok=True)
    shutil.copy2(output_file, output_dir / output_file.name)
    shutil.copy2(subject_mapping_md, output_dir / subject_mapping_md.name)
    shutil.copy2(subject_mapping_json, output_dir / subject_mapping_json.name)


def process_desensitization(
    *,
    task_id: str,
    input_path: Path,
    document_name: str,
    work_dir: Path,
    is_text: bool = False,
    output_dir: Path | None = None,
) -> dict[str, Path | dict[str, Any]]:
    work_dir.mkdir(parents=True, exist_ok=True)
    engine = Desensitizer()
    findings: list[Finding] = []
    subject_mappings: list[SubjectMapping] = []
    warnings: list[str] = []
    suffix = '.txt' if is_text else input_path.suffix.lower()
    input_name = input_path.name

    if is_text or suffix in TEXT_EXTENSIONS:
        raw = read_text_file(input_path)
        redacted, text_findings, text_subjects = sanitize_text_and_subjects(
            raw, engine, surface='text', locator='全文'
        )
        findings.extend(text_findings)
        subject_mappings.extend(text_subjects)
        output_file = work_dir / f'desensitized_output{text_output_suffix(input_path)}'
        output_file.write_text(redacted, encoding='utf-8')
        input_type = 'text'
    elif suffix in DOC_EXTENSIONS:
        output_file, doc_findings, doc_subjects = process_docx(input_path, work_dir, engine, warnings)
        findings.extend(doc_findings)
        subject_mappings.extend(doc_subjects)
        input_type = 'document'
    elif suffix in PDF_EXTENSIONS:
        output_file, pdf_findings, pdf_subjects = process_pdf(input_path, work_dir, engine, warnings)
        findings.extend(pdf_findings)
        subject_mappings.extend(pdf_subjects)
        input_type = 'pdf'
    elif suffix in TABLE_EXTENSIONS:
        output_file, table_findings, table_subjects = process_table(input_path, work_dir, engine, warnings)
        findings.extend(table_findings)
        subject_mappings.extend(table_subjects)
        input_type = 'table'
    elif suffix in JSON_EXTENSIONS:
        output_file, json_findings, json_subjects = process_json(input_path, work_dir, engine, warnings)
        findings.extend(json_findings)
        subject_mappings.extend(json_subjects)
        input_type = 'json'
    elif suffix in PRESENTATION_EXTENSIONS:
        output_file, presentation_findings, presentation_subjects = process_pptx(input_path, work_dir, engine, warnings)
        findings.extend(presentation_findings)
        subject_mappings.extend(presentation_subjects)
        input_type = 'presentation'
    elif suffix in IMAGE_EXTENSIONS:
        output_file, image_findings, image_subjects = process_image(input_path, work_dir, engine, warnings)
        findings.extend(image_findings)
        subject_mappings.extend(image_subjects)
        input_type = 'image'
    else:
        raise ValueError(f'暂不支持该文件类型：{suffix or "无扩展名"}')

    report = build_report(
        task_id=task_id,
        document_name=document_name,
        input_name=input_name,
        input_type=input_type,
        output_file=output_file,
        findings=findings,
        warnings=warnings,
        engine=engine,
    )
    report_json = work_dir / 'desensitization_report.json'
    report_md = work_dir / 'desensitization_report.md'
    note = work_dir / 'original_retention_note.txt'
    subject_mapping_md = work_dir / 'subject_mapping.md'
    subject_mapping_json = work_dir / 'subject_mapping.json'
    report_json.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding='utf-8')
    report_md.write_text(render_report_markdown(report), encoding='utf-8')
    write_retention_note(note)
    subject_mapping_md.write_text(
        render_subject_mapping_markdown(task_id, document_name, subject_mappings),
        encoding='utf-8',
    )
    subject_mapping_json.write_text(
        json.dumps(build_subject_mapping_json(task_id, document_name, subject_mappings), ensure_ascii=False, indent=2),
        encoding='utf-8',
    )

    if output_dir is not None:
        copy_outputs_to_directory(
            output_dir,
            output_file,
            subject_mapping_md,
            subject_mapping_json,
            document_name,
        )

    return {
        'output_file': output_file,
        'report_json': report_json,
        'report_md': report_md,
        'retention_note': note,
        'subject_mapping_md': subject_mapping_md,
        'subject_mapping_json': subject_mapping_json,
        'report': report,
    }


def process_docx(
    input_path: Path,
    work_dir: Path,
    engine: Desensitizer,
    warnings: list[str],
) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    all_findings: list[Finding] = []
    all_subjects: list[SubjectMapping] = []

    if input_path.suffix.lower() == '.doc':
        text = read_text_file(input_path)
        redacted, findings, subjects = sanitize_text_and_subjects(
            text, engine, surface='doc', locator='全文'
        )
        all_findings.extend(findings)
        all_subjects.extend(subjects)
        output_file = work_dir / 'desensitized_output.txt'
        output_file.write_text(redacted, encoding='utf-8')
        warnings.append('.doc 文件按纯文本兜底解析，复杂格式可能无法保留。')
        return output_file, all_findings, all_subjects

    document = Document(str(input_path))
    for index, paragraph in enumerate(document.paragraphs, start=1):
        if not paragraph.text:
            continue
        redacted, findings, subjects = sanitize_text_and_subjects(
            paragraph.text, engine, surface='docx', locator=f'段落 {index}'
        )
        all_findings.extend(findings)
        all_subjects.extend(subjects)
        if findings or subjects:
            paragraph.text = redacted

    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            for col_index, cell in enumerate(row.cells, start=1):
                if not cell.text:
                    continue
                locator = f'表格 {table_index} 行 {row_index} 列 {col_index}'
                redacted, findings, subjects = sanitize_text_and_subjects(
                    cell.text, engine, surface='docx_table', locator=locator
                )
                all_findings.extend(findings)
                all_subjects.extend(subjects)
                if findings or subjects:
                    cell.text = redacted

    output_file = work_dir / 'desensitized_output.docx'
    document.save(str(output_file))
    return output_file, all_findings, all_subjects


def _expanded_pdf_rect(rect: Any, page: Any, padding: float = REDACTION_PADDING) -> Any:
    bounds = page.rect
    expanded = fitz.Rect(rect)
    expanded.x0 = max(bounds.x0, expanded.x0 - padding)
    expanded.y0 = max(bounds.y0, expanded.y0 - padding)
    expanded.x1 = min(bounds.x1, expanded.x1 + padding)
    expanded.y1 = min(bounds.y1, expanded.y1 + padding)
    return expanded


def _search_pdf_rects(page: Any, value: str) -> list[Any]:
    if not value or not value.strip():
        return []
    candidates = [
        value,
        re.sub(r'\s+', ' ', value).strip(),
        re.sub(r'\s+', '', value).strip(),
    ]
    seen_candidates: set[str] = set()
    rects: list[Any] = []
    seen_rects: set[tuple[int, int, int, int]] = set()
    for candidate in candidates:
        if not candidate or candidate in seen_candidates:
            continue
        seen_candidates.add(candidate)
        try:
            found = page.search_for(candidate)
        except Exception:
            found = []
        for rect in found:
            key = (round(rect.x0), round(rect.y0), round(rect.x1), round(rect.y1))
            if key in seen_rects:
                continue
            seen_rects.add(key)
            rects.append(rect)
    return rects


def _add_pdf_text_redactions(page: Any, values: list[str]) -> int:
    count = 0
    for value in values:
        for rect in _search_pdf_rects(page, value):
            page.add_redact_annot(_expanded_pdf_rect(rect, page), fill=PDF_REDACTION_FILL)
            count += 1
    return count


def _add_pdf_branding_image_redactions(page: Any, page_index: int, page_text: str) -> int:
    """封面上方的小图片多为主体 logo，命中主体/Originator 语境时一并遮盖。"""
    if page_index > 2 or not re.search(r'originator|servicer|finance|company|co\.|ltd|公司', page_text, re.IGNORECASE):
        return 0
    if not hasattr(page, 'get_image_info'):
        return 0

    page_area = max(1.0, float(page.rect.width * page.rect.height))
    count = 0
    try:
        image_infos = page.get_image_info(xrefs=True)
    except Exception:
        return 0

    for info in image_infos:
        bbox = info.get('bbox')
        if not bbox:
            continue
        rect = fitz.Rect(bbox)
        if rect.is_empty:
            continue
        image_area = float(rect.width * rect.height)
        if rect.y0 > page.rect.height * 0.42:
            continue
        if image_area > page_area * 0.12:
            continue
        if rect.width < 8 or rect.height < 8:
            continue
        page.add_redact_annot(_expanded_pdf_rect(rect, page, padding=2.0), fill=PDF_REDACTION_FILL)
        count += 1
    return count


def _apply_pdf_redactions(page: Any) -> None:
    image_mode = getattr(fitz, 'PDF_REDACT_IMAGE_PIXELS', None)
    try:
        if image_mode is not None:
            page.apply_redactions(images=image_mode)
        else:
            page.apply_redactions()
    except TypeError:
        page.apply_redactions()


def process_pdf(
    input_path: Path,
    work_dir: Path,
    engine: Desensitizer,
    warnings: list[str],
) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    text_parts: list[str] = []
    all_findings: list[Finding] = []
    all_subjects: list[SubjectMapping] = []
    if fitz is not None:
        pdf = fitz.open(str(input_path))
        redaction_count = 0
        for page_index, page in enumerate(pdf, start=1):
            page_text = page.get_text('text') or ''
            if page_text.strip():
                locator = f'第 {page_index} 页'
                findings = engine._detect(page_text, surface='pdf', locator=locator)
                _redacted_text, subjects = redact_legal_subjects(
                    page_text,
                    locator=locator,
                    detector=engine._subject_detector,
                )
                all_findings.extend(findings)
                all_subjects.extend(subjects)
                subject_entities = detect_legal_subject_entities(page_text, detector=engine._subject_detector)
                values = [page_text[item.start:item.end] for item in findings]
                values.extend(str(getattr(entity, 'text', '')) for entity in subject_entities)
                for match in BRANDING_TEXT_PATTERN.finditer(page_text):
                    values.append(match.group(0))
                page_redactions = _add_pdf_text_redactions(page, values)
                page_redactions += _add_pdf_branding_image_redactions(page, page_index, page_text)
                redaction_count += page_redactions
                if page_redactions:
                    _apply_pdf_redactions(page)
                redacted_preview = replace_spans(page_text, findings)
                redacted_preview, _preview_subjects = redact_legal_subjects(
                    redacted_preview,
                    locator=locator,
                    detector=engine._subject_detector,
                )
                text_parts.append(f'--- 第 {page_index} 页 ---\n{redacted_preview}')

    if text_parts:
        output_file = work_dir / 'desensitized_output.pdf'
        try:
            pdf.save(str(output_file), garbage=4, deflate=True)
        finally:
            pdf.close()
        preview_file = work_dir / 'desensitized_output_preview.txt'
        preview_file.write_text('\n\n'.join(text_parts), encoding='utf-8')
        if redaction_count:
            warnings.append(f'PDF 已执行版式级遮盖 {redaction_count} 处；请抽样复核 logo 和扫描图像区域。')
        else:
            warnings.append('PDF 可复制文本未定位到可遮盖坐标，已保留脱敏文本预览供复核。')
        return output_file, all_findings, all_subjects
    if fitz is not None:
        pdf.close()

    warnings.append('PDF 未提取到可复制文本，尝试按扫描件图片进行 OCR 脱敏。')
    if fitz is None:
        raise RuntimeError('扫描型 PDF 需要 PyMuPDF 才能渲染页面，请先安装 PyMuPDF。')
    if Image is None or ImageDraw is None:
        raise RuntimeError('扫描型 PDF 或图片脱敏需要安装 Pillow。')
    try:
        ensure_ocr_available(require_pdf=True)
    except OcrUnavailable as exc:
        raise RuntimeError(str(exc)) from exc

    pdf = fitz.open(str(input_path))
    page_outputs: list[str] = []
    preview_parts: list[str] = []
    for page_index, page in enumerate(pdf, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        page_image = work_dir / f'_pdf_page_{page_index}.png'
        pix.save(str(page_image))
        page_result = redact_image_with_preview(
            page_image, work_dir / f'desensitized_output_page_{page_index}.png', engine, f'第 {page_index} 页'
        )
        page_outputs.append(page_result.output_file.name)
        if page_result.preview_text:
            preview_parts.append(f'--- 第 {page_index} 页 ---\n{page_result.preview_text}')
        all_findings.extend(page_result.findings)
        all_subjects.extend(page_result.subjects)
    pdf.close()
    preview_file = work_dir / 'desensitized_output_preview.txt'
    preview_file.write_text('\n\n'.join(preview_parts).rstrip() + ('\n' if preview_parts else ''), encoding='utf-8')
    manifest = work_dir / 'desensitized_output_pages.json'
    manifest.write_text(
        json.dumps(
            {
                'pages': page_outputs,
                'preview_text_file': preview_file.name,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding='utf-8',
    )
    return manifest, all_findings, all_subjects


def process_table(
    input_path: Path,
    work_dir: Path,
    engine: Desensitizer,
    warnings: list[str],
) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    if input_path.suffix.lower() == '.csv':
        return process_csv(input_path, work_dir, engine, delimiter=',', output_suffix='.csv')
    if input_path.suffix.lower() == '.tsv':
        return process_csv(input_path, work_dir, engine, delimiter='\t', output_suffix='.tsv')
    if pd is None:
        raise RuntimeError('处理 xlsx/xls/ods 需要安装 pandas、openpyxl、xlrd 或 odfpy。')

    sheets = pd.read_excel(input_path, sheet_name=None, dtype=str)
    output_file = work_dir / 'desensitized_output.xlsx'
    all_findings: list[Finding] = []
    all_subjects: list[SubjectMapping] = []
    with pd.ExcelWriter(output_file, engine='openpyxl') as writer:
        for sheet_name, frame in sheets.items():
            sanitized_frame = frame.copy()
            for row_index, row in sanitized_frame.iterrows():
                for column in sanitized_frame.columns:
                    value = '' if row[column] is None else str(row[column])
                    if not value or value == 'nan':
                        continue
                    locator = f'{sheet_name}!{column}{row_index + 2}'
                    redacted, findings, subjects = sanitize_text_and_subjects(
                        value, engine, surface='xlsx_cell', locator=locator
                    )
                    if findings or subjects:
                        sanitized_frame.at[row_index, column] = redacted
                        all_findings.extend(findings)
                        all_subjects.extend(subjects)
            sanitized_frame.to_excel(writer, sheet_name=str(sheet_name)[:31], index=False)
    return output_file, all_findings, all_subjects


def process_csv(
    input_path: Path,
    work_dir: Path,
    engine: Desensitizer,
    *,
    delimiter: str,
    output_suffix: str,
) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    output_file = work_dir / f'desensitized_output{output_suffix}'
    all_findings: list[Finding] = []
    all_subjects: list[SubjectMapping] = []
    with input_path.open('r', encoding='utf-8-sig', newline='') as source:
        reader = csv.reader(source, delimiter=delimiter)
        rows = list(reader)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row, start=1):
            redacted, findings, subjects = sanitize_text_and_subjects(
                value, engine, surface='csv_cell', locator=f'R{row_index}C{col_index}'
            )
            if findings or subjects:
                row[col_index - 1] = redacted
                all_findings.extend(findings)
                all_subjects.extend(subjects)
    with output_file.open('w', encoding='utf-8-sig', newline='') as target:
        writer = csv.writer(target, delimiter=delimiter)
        writer.writerows(rows)
    return output_file, all_findings, all_subjects


def process_json(input_path: Path, work_dir: Path, engine: Desensitizer, warnings: list[str]) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    if input_path.suffix.lower() in {'.jsonl', '.ndjson'}:
        return process_json_lines(input_path, work_dir, engine, warnings)

    payload = json.loads(input_path.read_text(encoding='utf-8'))
    all_findings: list[Finding] = []
    all_subjects: list[SubjectMapping] = []

    def walk(value: Any, path: str) -> Any:
        if isinstance(value, str):
            redacted, findings, subjects = sanitize_text_and_subjects(
                value, engine, surface='json_value', locator=path
            )
            all_findings.extend(findings)
            all_subjects.extend(subjects)
            return redacted
        if isinstance(value, list):
            return [walk(item, f'{path}[{index}]') for index, item in enumerate(value)]
        if isinstance(value, dict):
            return {key: walk(item, f'{path}.{key}' if path else str(key)) for key, item in value.items()}
        return value

    sanitized_payload = walk(payload, '')
    output_file = work_dir / 'desensitized_output.json'
    output_file.write_text(json.dumps(sanitized_payload, ensure_ascii=False, indent=2), encoding='utf-8')
    return output_file, all_findings, all_subjects


def process_json_lines(
    input_path: Path,
    work_dir: Path,
    engine: Desensitizer,
    warnings: list[str],
) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    all_findings: list[Finding] = []
    all_subjects: list[SubjectMapping] = []
    output_lines: list[str] = []

    def walk(value: Any, path: str) -> Any:
        if isinstance(value, str):
            redacted, findings, subjects = sanitize_text_and_subjects(
                value, engine, surface='jsonl_value', locator=path
            )
            all_findings.extend(findings)
            all_subjects.extend(subjects)
            return redacted
        if isinstance(value, list):
            return [walk(item, f'{path}[{index}]') for index, item in enumerate(value)]
        if isinstance(value, dict):
            return {key: walk(item, f'{path}.{key}' if path else str(key)) for key, item in value.items()}
        return value

    for line_index, line in enumerate(input_path.read_text(encoding='utf-8').splitlines(), start=1):
        if not line.strip():
            output_lines.append(line)
            continue
        try:
            payload = json.loads(line)
        except json.JSONDecodeError:
            redacted, findings, subjects = sanitize_text_and_subjects(
                line, engine, surface='jsonl_line', locator=f'行 {line_index}'
            )
            all_findings.extend(findings)
            all_subjects.extend(subjects)
            output_lines.append(redacted)
            warnings.append(f'第 {line_index} 行不是合法 JSON，已按普通文本脱敏。')
            continue
        output_lines.append(json.dumps(walk(payload, f'行 {line_index}'), ensure_ascii=False))

    output_file = work_dir / f'desensitized_output{input_path.suffix.lower()}'
    output_file.write_text('\n'.join(output_lines) + ('\n' if output_lines else ''), encoding='utf-8')
    return output_file, all_findings, all_subjects


def process_pptx(
    input_path: Path,
    work_dir: Path,
    engine: Desensitizer,
    warnings: list[str],
) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    if Presentation is None:
        raise RuntimeError('处理 pptx 需要安装 python-pptx。')

    presentation = Presentation(str(input_path))
    all_findings: list[Finding] = []
    all_subjects: list[SubjectMapping] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, 'has_text_frame', False) and shape.text:
                locator = f'幻灯片 {slide_index} 文本框 {shape_index}'
                redacted, findings, subjects = sanitize_text_and_subjects(
                    shape.text, engine, surface='pptx_text', locator=locator
                )
                if findings or subjects:
                    shape.text = redacted
                    all_findings.extend(findings)
                    all_subjects.extend(subjects)
            if getattr(shape, 'has_table', False):
                for row_index, row in enumerate(shape.table.rows, start=1):
                    for col_index, cell in enumerate(row.cells, start=1):
                        if not cell.text:
                            continue
                        locator = f'幻灯片 {slide_index} 表格 {shape_index} 行 {row_index} 列 {col_index}'
                        redacted, findings, subjects = sanitize_text_and_subjects(
                            cell.text, engine, surface='pptx_table', locator=locator
                        )
                        if findings or subjects:
                            cell.text = redacted
                            all_findings.extend(findings)
                            all_subjects.extend(subjects)

    output_file = work_dir / 'desensitized_output.pptx'
    presentation.save(str(output_file))
    warnings.append('PPTX 脱敏会尽量保留结构，但命中文本框的局部字体样式可能被重置。')
    return output_file, all_findings, all_subjects


def process_image(
    input_path: Path,
    work_dir: Path,
    engine: Desensitizer,
    warnings: list[str],
) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    if Image is None or ImageDraw is None:
        raise RuntimeError('图片脱敏需要安装 Pillow。')
    try:
        ensure_ocr_available()
    except OcrUnavailable as exc:
        raise RuntimeError(str(exc)) from exc
    output_file = work_dir / f'desensitized_output{input_path.suffix.lower()}'
    return redact_image(input_path, output_file, engine, '图片')


def _ocr_token_rows(data: dict[str, Any]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    texts = data.get('text', [])
    for index, raw in enumerate(texts):
        text = (raw or '').strip()
        if not text:
            continue
        try:
            left = int(float(data.get('left', [])[index]))
            top = int(float(data.get('top', [])[index]))
            width = int(float(data.get('width', [])[index]))
            height = int(float(data.get('height', [])[index]))
        except Exception:
            continue
        rows.append({
            'index': index,
            'text': text,
            'left': left,
            'top': top,
            'width': max(1, width),
            'height': max(1, height),
            'block_num': data.get('block_num', [0] * len(texts))[index] if index < len(data.get('block_num', [])) else 0,
            'par_num': data.get('par_num', [0] * len(texts))[index] if index < len(data.get('par_num', [])) else 0,
            'line_num': data.get('line_num', [index] * len(texts))[index] if index < len(data.get('line_num', [])) else index,
            'word_num': data.get('word_num', [index] * len(texts))[index] if index < len(data.get('word_num', [])) else index,
        })
    return rows


def _group_ocr_lines(tokens: list[dict[str, Any]]) -> list[dict[str, Any]]:
    groups: dict[tuple[Any, Any, Any], list[dict[str, Any]]] = defaultdict(list)
    for token in tokens:
        key = (token['block_num'], token['par_num'], token['line_num'])
        groups[key].append(token)

    lines: list[dict[str, Any]] = []
    for key, items in groups.items():
        ordered = sorted(items, key=lambda item: (item['top'], item['left'], item['word_num']))
        parts: list[str] = []
        spans: list[tuple[int, int, dict[str, Any]]] = []
        cursor = 0
        for item in ordered:
            if parts and _needs_ocr_token_space(parts[-1], item['text']):
                parts.append(' ')
                cursor += 1
            start = cursor
            parts.append(item['text'])
            cursor += len(item['text'])
            spans.append((start, cursor, item))
        lines.append({'key': key, 'text': ''.join(parts), 'spans': spans, 'tokens': ordered})
    return sorted(lines, key=lambda line: (min(t['top'] for t in line['tokens']), min(t['left'] for t in line['tokens'])))


def _is_cjk_or_punctuation(char: str) -> bool:
    return bool(re.match(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef，。；：？！、）】》」』]', char))


def _needs_ocr_token_space(previous: str, current: str) -> bool:
    if not previous or not current:
        return False
    if _is_cjk_or_punctuation(previous[-1]) or _is_cjk_or_punctuation(current[0]):
        return False
    return True


def _redact_ocr_line_preview(
    text: str,
    findings: list[Finding],
    engine: Desensitizer,
    locator: str,
) -> str:
    redacted = replace_spans(text, findings)
    redacted, _subjects = redact_legal_subjects(
        redacted,
        locator=locator,
        detector=engine._subject_detector,
    )
    return redacted


def _span_intersects(start: int, end: int, token_start: int, token_end: int) -> bool:
    return start < token_end and end > token_start


def _draw_ocr_span(
    draw: Any,
    line: dict[str, Any],
    start: int,
    end: int,
    image_width: int,
    image_height: int,
) -> bool:
    spans = line['spans']
    if not spans:
        return False

    if len(spans) == 1:
        token = spans[0][2]
        text_len = max(1, len(line['text']))
        ratio_start = max(0.0, min(1.0, start / text_len))
        ratio_end = max(ratio_start + 0.02, min(1.0, end / text_len))
        left = token['left'] + int(token['width'] * ratio_start)
        right = token['left'] + int(token['width'] * ratio_end)
        top = token['top']
        bottom = token['top'] + token['height']
        draw.rectangle([
            max(0, left - 2),
            max(0, top - 2),
            min(image_width, right + 2),
            min(image_height, bottom + 2),
        ], fill=REDACTION_FILL)
        return True

    drawn = False
    for token_start, token_end, token in spans:
        if not _span_intersects(start, end, token_start, token_end):
            continue
        draw.rectangle([
            max(0, token['left'] - 2),
            max(0, token['top'] - 2),
            min(image_width, token['left'] + token['width'] + 2),
            min(image_height, token['top'] + token['height'] + 2),
        ], fill=REDACTION_FILL)
        drawn = True
    return drawn


def redact_image_with_preview(
    input_path: Path,
    output_file: Path,
    engine: Desensitizer,
    locator_prefix: str,
) -> ImageRedactionResult:
    image = Image.open(input_path).convert('RGB')
    data = ocr_image_to_data(image)
    draw = ImageDraw.Draw(image)
    all_findings: list[Finding] = []
    all_subjects: list[SubjectMapping] = []
    preview_lines: list[str] = []
    tokens = _ocr_token_rows(data)
    lines = _group_ocr_lines(tokens)
    for line_index, line in enumerate(lines, start=1):
        text = line['text']
        locator = f'{locator_prefix} OCR行#{line_index}'
        _redacted, findings, subjects = sanitize_text_and_subjects(
            text, engine, surface='image_text', locator=locator
        )
        all_findings.extend(findings)
        all_subjects.extend(subjects)
        preview_lines.append(_redact_ocr_line_preview(text, findings, engine, locator))
        subject_entities = detect_legal_subject_entities(text, detector=engine._subject_detector)
        spans_to_redact: list[tuple[int, int]] = [(item.start, item.end) for item in findings]
        spans_to_redact.extend((int(getattr(entity, 'start', 0)), int(getattr(entity, 'end', 0))) for entity in subject_entities)
        spans_to_redact.extend((match.start(), match.end()) for match in BRANDING_TEXT_PATTERN.finditer(text))
        for start, end in spans_to_redact:
            if end <= start:
                continue
            _draw_ocr_span(draw, line, start, end, image.width, image.height)
    image.save(output_file)
    return ImageRedactionResult(
        output_file=output_file,
        findings=all_findings,
        subjects=all_subjects,
        preview_text='\n'.join(line for line in preview_lines if line.strip()).strip(),
    )


def redact_image(input_path: Path, output_file: Path, engine: Desensitizer, locator_prefix: str) -> tuple[Path, list[Finding], list[SubjectMapping]]:
    result = redact_image_with_preview(input_path, output_file, engine, locator_prefix)
    return result.output_file, result.findings, result.subjects
