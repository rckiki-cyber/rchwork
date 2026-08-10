#!/usr/bin/env python3
import importlib.util
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation


SCRIPT_DIR = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))
SCRIPT = SCRIPT_DIR / "local_export_pptx.py"
SPEC = importlib.util.spec_from_file_location("local_export_pptx_under_test", SCRIPT)
MODULE = importlib.util.module_from_spec(SPEC)
assert SPEC and SPEC.loader
SPEC.loader.exec_module(MODULE)


class LocalExportPptxTests(unittest.TestCase):
    def create_project(self, root: Path) -> Path:
        project = root / "deck"
        (project / "pages").mkdir(parents=True)
        (project / "deck.pptd").write_text(
            """version: v2
title: Local deck
size: [960, 540]
theme:
  colors: {primary: "#172554", accent: "#EA580C", text: "#111827"}
  textStyles:
    title: {fontSize: 38, color: "$primary", fontFamily: Arial, bold: true}
    body: {fontSize: 18, color: "$text", fontFamily: Arial}
  tableStyles:
    default:
      firstRowStyle: {fill: {type: solid, color: "$primary"}, color: "#FFFFFF", bold: true}
pages: [pages/01.page]
""",
            encoding="utf-8",
        )
        (project / "pages" / "01.page").write_text(
            """pageType: content
background: {type: solid, color: "#F8FAFC"}
elements:
  - elementId: heading
    elementType: text
    bounds: [60, 35, 840, 65]
    content:
      style: "$title"
      align: [center, middle]
      text: '<p><strong>本地 Kimi 风格导出</strong></p>'
  - elementId: card
    elementType: shape
    bounds: [60, 125, 390, 145]
    shapeName: roundRect
    fill: {type: solid, color: "$primary"}
    border: {style: solid, width: 2, color: "$accent"}
  - elementId: connector
    elementType: line
    bounds: [450, 185, 60, 20]
    viewBox: [1, 1]
    points: "0,0.5 1,0.5"
    border: {style: dash, width: 2, color: "$accent"}
  - elementId: metrics
    elementType: table
    bounds: [510, 125, 390, 145]
    columnWidths: [0.5, 0.5]
    rowHeights: [0.5, 0.5]
    style: "$default"
    rows:
      - [{text: 指标}, {text: 数值}]
      - [{text: 完成率}, {text: 96%}]
  - elementId: trend
    elementType: chart
    bounds: [180, 305, 600, 190]
    data:
      cols: [季度, 实际]
      rows: [[Q1, 12], [Q2, 18], [Q3, 25]]
    series:
      - {type: bar, encode: {x: 季度, y: 实际}, name: 实际, fill: "$accent"}
""",
            encoding="utf-8",
        )
        return project / "deck.pptd"

    def test_local_export_preserves_pptd_structure_and_verifies_zip(self):
        with tempfile.TemporaryDirectory() as name:
            manifest = self.create_project(Path(name))
            output = manifest.with_suffix(".pptx")
            result = MODULE.export_local_pptx(manifest, output, "fade")
            self.assertEqual(result["exporter"], "local-python-pptx")
            self.assertEqual(result["slides"], 1)
            self.assertEqual(result["fadeTransitions"], 1)
            self.assertGreater(result["bytes"], 10_000)
            self.assertFalse(any("本地渲染失败" in warning for warning in result["warnings"]))
            deck = Presentation(str(output))
            self.assertEqual(len(deck.slides), 1)
            self.assertGreaterEqual(len(deck.slides[0].shapes), 5)
            text = "\n".join(
                shape.text for shape in deck.slides[0].shapes
                if getattr(shape, "has_text_frame", False)
            )
            self.assertIn("本地 Kimi 风格导出", text)
            table_text = "\n".join(
                cell.text
                for shape in deck.slides[0].shapes if getattr(shape, "has_table", False)
                for row in shape.table.rows for cell in row.cells
            )
            self.assertIn("完成率", table_text)
            with zipfile.ZipFile(output) as archive:
                self.assertIsNone(archive.testzip())
                self.assertIn(b"<p:fade/>", archive.read("ppt/slides/slide1.xml"))

    def test_existing_output_is_not_overwritten_without_force(self):
        with tempfile.TemporaryDirectory() as name:
            manifest = self.create_project(Path(name))
            output = manifest.with_suffix(".pptx")
            output.write_bytes(b"keep")
            with self.assertRaisesRegex(MODULE.ExportError, "already exists"):
                MODULE.export_local_pptx(manifest, output, "none")
            self.assertEqual(output.read_bytes(), b"keep")


if __name__ == "__main__":
    unittest.main()
