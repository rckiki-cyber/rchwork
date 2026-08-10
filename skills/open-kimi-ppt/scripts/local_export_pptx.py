#!/usr/bin/env python3
"""Deterministic local PPTD v2 to PPTX exporter.

This is the delivery fallback for the unified open-kimi-ppt Skill. It consumes
the exact same positioned PPTD project as the Kimi browser writer; it is not a
generic title/bullets template. Unsupported browser-only effects are reported
as warnings instead of silently dropping the affected element.
"""

from __future__ import annotations

import html
import re
import tempfile
import urllib.parse
import urllib.request
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

PPTX_IMPORT_ERROR: Optional[ImportError] = None
try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - exercised only on broken installs
    PPTX_IMPORT_ERROR = exc

from export_pptx import ExportError, patch_transitions, read_yaml_mapping, verify_output


DEFAULT_SLIDE_WIDTH_IN = 13.333333
CSS_STYLE_RE = re.compile(r"\s*([^:;]+)\s*:\s*([^;]+)")
POINT_RE = re.compile(r"-?\d+(?:\.\d+)?")


def _as_dict(value: Any) -> Dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _resolve_ref(value: Any, mapping: Dict[str, Any]) -> Any:
    if isinstance(value, str) and value.startswith("$"):
        return mapping.get(value[1:], value)
    return value


def _hex_color(value: Any, colors: Dict[str, Any], default: str = "#000000") -> RGBColor:
    resolved = str(_resolve_ref(value, colors) or default).strip()
    match = re.fullmatch(r"#?([0-9a-fA-F]{6})(?:[0-9a-fA-F]{2})?", resolved)
    if not match:
        match = re.fullmatch(r"#?([0-9a-fA-F]{3})", resolved)
        if match:
            short = match.group(1)
            return RGBColor.from_string("".join(ch * 2 for ch in short).upper())
        return RGBColor.from_string(default.lstrip("#")[:6].upper())
    return RGBColor.from_string(match.group(1).upper())


def _parse_css(style: str) -> Dict[str, Any]:
    result: Dict[str, Any] = {}
    for item in style.split(";"):
        match = CSS_STYLE_RE.fullmatch(item)
        if not match:
            continue
        key, value = match.group(1).strip().lower(), match.group(2).strip()
        if key == "font-size":
            number = POINT_RE.search(value)
            if number:
                result["fontSize"] = float(number.group())
        elif key == "font-family":
            result["fontFamily"] = value.strip("'\"")
        elif key == "color":
            result["color"] = value
        elif key == "font-weight":
            result["bold"] = value.lower() in {"bold", "bolder"} or value.isdigit() and int(value) >= 600
        elif key == "font-style":
            result["italic"] = value.lower() == "italic"
        elif key == "text-align":
            result["paragraphAlign"] = value.lower()
    return result


class _RichTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.paragraphs: List[Dict[str, Any]] = []
        self.current: Optional[Dict[str, Any]] = None
        self.styles: List[Dict[str, Any]] = [{}]

    def _paragraph(self) -> Dict[str, Any]:
        if self.current is None:
            self.current = {"runs": [], "style": {}}
            self.paragraphs.append(self.current)
        return self.current

    def handle_starttag(self, tag: str, attrs: List[Tuple[str, Optional[str]]]) -> None:
        attr = dict(attrs)
        if tag == "p":
            self.current = {"runs": [], "style": _parse_css(attr.get("style") or "")}
            self.paragraphs.append(self.current)
        elif tag == "br":
            self._paragraph()["runs"].append({"text": "\n", "style": dict(self.styles[-1])})
        if tag in {"span", "strong", "b", "em", "i"}:
            style = dict(self.styles[-1])
            style.update(_parse_css(attr.get("style") or ""))
            if tag in {"strong", "b"}:
                style["bold"] = True
            if tag in {"em", "i"}:
                style["italic"] = True
            self.styles.append(style)

    def handle_endtag(self, tag: str) -> None:
        if tag == "p":
            self.current = None
        if tag in {"span", "strong", "b", "em", "i"} and len(self.styles) > 1:
            self.styles.pop()

    def handle_data(self, data: str) -> None:
        if data and (self.current is not None or data.strip()):
            self._paragraph()["runs"].append({"text": data, "style": dict(self.styles[-1])})


def _rich_paragraphs(value: Any) -> List[Dict[str, Any]]:
    text = "" if value is None else str(value)
    if "<" not in text:
        return [
            {"runs": [{"text": part, "style": {}}], "style": {}}
            for part in text.splitlines() or [""]
        ]
    parser = _RichTextParser()
    try:
        parser.feed(text)
        parser.close()
    except Exception:
        return [{"runs": [{"text": html.unescape(re.sub(r"<[^>]+>", "", text)), "style": {}}], "style": {}}]
    return parser.paragraphs or [{"runs": [{"text": "", "style": {}}], "style": {}}]


class LocalPptdExporter:
    def __init__(self, manifest: Path) -> None:
        self.manifest = manifest.resolve()
        _text, self.deck = read_yaml_mapping(self.manifest)
        self.root = self.manifest.parent
        self.theme = _as_dict(self.deck.get("theme"))
        self.colors = _as_dict(self.theme.get("colors"))
        self.text_styles = _as_dict(self.theme.get("textStyles"))
        self.table_styles = _as_dict(self.theme.get("tableStyles"))
        size = self.deck.get("size")
        if not isinstance(size, list) or len(size) != 2:
            raise ExportError("PPTD size must contain [width, height]")
        self.canvas_width = float(size[0])
        self.canvas_height = float(size[1])
        self.slide_width_in = DEFAULT_SLIDE_WIDTH_IN
        self.slide_height_in = self.slide_width_in * self.canvas_height / self.canvas_width
        self.warnings: List[str] = []
        self._temp_files: List[Path] = []

    def emu_x(self, value: float) -> int:
        return int(Inches(float(value) / self.canvas_width * self.slide_width_in))

    def emu_y(self, value: float) -> int:
        return int(Inches(float(value) / self.canvas_height * self.slide_height_in))

    def rect(self, bounds: Sequence[float]) -> Tuple[int, int, int, int]:
        return self.emu_x(bounds[0]), self.emu_y(bounds[1]), self.emu_x(bounds[2]), self.emu_y(bounds[3])

    def style(self, content: Dict[str, Any]) -> Dict[str, Any]:
        result: Dict[str, Any] = {}
        reference = content.get("style") or content.get("textStyle")
        resolved = _resolve_ref(reference, self.text_styles)
        if isinstance(resolved, dict):
            result.update(resolved)
        for key in (
            "color", "fontSize", "fontFamily", "bold", "italic", "backgroundColor",
            "lineHeight", "lineHeightPx", "letterSpacing", "marginTop",
        ):
            if key in content:
                result[key] = content[key]
        return result

    def apply_font(self, font: Any, style: Dict[str, Any]) -> None:
        font.name = str(style.get("fontFamily") or "Microsoft YaHei")
        try:
            font.size = Pt(float(style.get("fontSize", 18)))
        except (TypeError, ValueError):
            font.size = Pt(18)
        font.bold = bool(style.get("bold", False))
        font.italic = bool(style.get("italic", False))
        font.color.rgb = _hex_color(style.get("color", "$text"), self.colors, "#111827")

    def apply_fill(self, target: Any, fill_spec: Any, label: str) -> None:
        fill = _as_dict(fill_spec)
        if not fill:
            target.fill.background()
            return
        fill_type = fill.get("type", "solid")
        if fill_type == "solid":
            target.fill.solid()
            target.fill.fore_color.rgb = _hex_color(fill.get("color"), self.colors, "#FFFFFF")
            return
        if fill_type == "gradient":
            stops = fill.get("stops") if isinstance(fill.get("stops"), list) else []
            first = _as_dict(stops[0]).get("color") if stops else "#FFFFFF"
            target.fill.solid()
            target.fill.fore_color.rgb = _hex_color(first, self.colors, "#FFFFFF")
            self.warnings.append(f"{label}: gradient approximated with its first color")
            return
        target.fill.background()
        self.warnings.append(f"{label}: image fill is not supported locally; transparent fill used")

    def apply_border(self, line: Any, border_spec: Any) -> None:
        border = _as_dict(border_spec)
        if not border:
            line.fill.background()
            return
        line.color.rgb = _hex_color(border.get("color"), self.colors, "#000000")
        try:
            line.width = Pt(float(border.get("width", 1)))
        except (TypeError, ValueError):
            line.width = Pt(1)
        dash = str(border.get("style", "solid"))
        if dash == "dash":
            line.dash_style = MSO_LINE_DASH_STYLE.DASH
        elif dash == "dot":
            line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

    def add_text(self, slide: Any, element: Dict[str, Any]) -> None:
        x, y, w, h = self.rect(element["bounds"])
        shape = slide.shapes.add_textbox(x, y, w, h)
        shape.rotation = float(element.get("rotation", 0) or 0)
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = bool(_as_dict(element.get("content")).get("wrap", True))
        content = _as_dict(element.get("content"))
        align = content.get("align") if isinstance(content.get("align"), list) else ["left", "top"]
        frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM,
        }.get(str(align[1] if len(align) > 1 else "top"), MSO_ANCHOR.TOP)
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        base = self.style(content)
        paragraphs = _rich_paragraphs(content.get("text"))
        for index, record in enumerate(paragraphs):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            p_align = _as_dict(record.get("style")).get("paragraphAlign") or (align[0] if align else "left")
            paragraph.alignment = {
                "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }.get(str(p_align), PP_ALIGN.LEFT)
            line_height = base.get("lineHeight")
            if isinstance(line_height, (int, float)) and line_height > 0:
                paragraph.line_spacing = float(line_height)
            for run_record in record.get("runs", []):
                run = paragraph.add_run()
                run.text = str(run_record.get("text", ""))
                merged = dict(base)
                merged.update(_as_dict(run_record.get("style")))
                self.apply_font(run.font, merged)

    def add_shape(self, slide: Any, element: Dict[str, Any]) -> None:
        shape_map = {
            "rect": MSO_SHAPE.RECTANGLE, "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
            "ellipse": MSO_SHAPE.OVAL, "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND, "homePlate": MSO_SHAPE.PENTAGON,
            "chevron": MSO_SHAPE.CHEVRON, "donut": MSO_SHAPE.DONUT,
            "star5": MSO_SHAPE.STAR_5_POINT, "rightArrow": MSO_SHAPE.RIGHT_ARROW,
            "wedgeRectCallout": MSO_SHAPE.RECTANGULAR_CALLOUT,
            "bracePair": MSO_SHAPE.DOUBLE_BRACE,
        }
        name = str(element.get("shapeName", "rect"))
        if name not in shape_map:
            self.warnings.append(f"{element.get('elementId')}: shape {name!r} approximated as rectangle")
        x, y, w, h = self.rect(element["bounds"])
        shape = slide.shapes.add_shape(shape_map.get(name, MSO_SHAPE.RECTANGLE), x, y, w, h)
        shape.rotation = float(element.get("rotation", 0) or 0)
        self.apply_fill(shape, element.get("fill"), str(element.get("elementId")))
        self.apply_border(shape.line, element.get("border"))

    def add_line(self, slide: Any, element: Dict[str, Any]) -> None:
        x, y, w, h = self.rect(element["bounds"])
        points = re.findall(r"(-?\d+(?:\.\d+)?),(-?\d+(?:\.\d+)?)", str(element.get("points", "")))
        view = element.get("viewBox") if isinstance(element.get("viewBox"), list) else [1, 1]
        if len(points) >= 2 and float(view[0]) and float(view[1]):
            first, last = points[0], points[-1]
            x1 = x + int(w * float(first[0]) / float(view[0]))
            y1 = y + int(h * float(first[1]) / float(view[1]))
            x2 = x + int(w * float(last[0]) / float(view[0]))
            y2 = y + int(h * float(last[1]) / float(view[1]))
        else:
            x1, y1, x2, y2 = x, y, x + w, y + h
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        self.apply_border(connector.line, element.get("border") or {"color": "#000000", "width": 1})
        if len(points) > 2:
            self.warnings.append(f"{element.get('elementId')}: bezier line approximated as a straight connector")

    def resolve_media(self, src: Any) -> Optional[Path]:
        if not isinstance(src, str) or not src:
            return None
        if re.match(r"https?://", src):
            try:
                suffix = Path(urllib.parse.urlparse(src).path).suffix or ".img"
            except ValueError:
                suffix = Path(src.split("?", 1)[0]).suffix or ".img"
            try:
                handle = tempfile.NamedTemporaryFile(prefix="open-kimi-local-", suffix=suffix, delete=False)
                handle.close()
                target = Path(handle.name)
                urllib.request.urlretrieve(src, target)
                self._temp_files.append(target)
                return target
            except OSError as exc:
                self.warnings.append(f"remote image could not be downloaded: {src} ({exc})")
                return None
        path = (self.root / src).resolve()
        try:
            path.relative_to(self.root)
        except ValueError:
            self.warnings.append(f"image path escapes project and was skipped: {src}")
            return None
        if not path.is_file():
            self.warnings.append(f"image file is missing: {src}")
            return None
        return path

    def add_image(self, slide: Any, element: Dict[str, Any]) -> None:
        media = self.resolve_media(element.get("src"))
        if media is None:
            self.add_missing_element(slide, element, "图片不可用")
            return
        x, y, w, h = self.rect(element["bounds"])
        picture = slide.shapes.add_picture(str(media), x, y, w, h)
        picture.rotation = float(element.get("rotation", 0) or 0)
        crop = _as_dict(element.get("crop"))
        picture.crop_left = float(crop.get("left", 0) or 0)
        picture.crop_top = float(crop.get("top", 0) or 0)
        picture.crop_right = float(crop.get("right", 0) or 0)
        picture.crop_bottom = float(crop.get("bottom", 0) or 0)
        fit = _as_dict(element.get("fit")).get("mode", "cover")
        if fit != "fill":
            self.warnings.append(f"{element.get('elementId')}: image fit={fit} approximated to the PPTD bounds")

    def table_style(self, element: Dict[str, Any]) -> Dict[str, Any]:
        style = element.get("style")
        resolved = _resolve_ref(style, self.table_styles)
        return resolved if isinstance(resolved, dict) else {}

    def cell_style(self, table_style: Dict[str, Any], cell: Dict[str, Any], row: int, rows: int) -> Dict[str, Any]:
        result = dict(_as_dict(table_style.get("cellStyle")))
        body = table_style.get("bodyStyles")
        if row > 0 and isinstance(body, list) and body:
            result.update(_as_dict(body[(row - 1) % len(body)]))
        if row == 0:
            result.update(_as_dict(table_style.get("firstRowStyle")))
        if row == rows - 1:
            result.update(_as_dict(table_style.get("lastRowStyle")))
        result.update(cell)
        return result

    def add_table(self, slide: Any, element: Dict[str, Any]) -> None:
        rows_data = element.get("rows") if isinstance(element.get("rows"), list) else []
        column_widths = element.get("columnWidths") if isinstance(element.get("columnWidths"), list) else []
        row_heights = element.get("rowHeights") if isinstance(element.get("rowHeights"), list) else []
        row_count, col_count = len(row_heights) or len(rows_data), len(column_widths)
        if row_count <= 0 or col_count <= 0:
            self.add_missing_element(slide, element, "表格数据为空")
            return
        x, y, w, h = self.rect(element["bounds"])
        table = slide.shapes.add_table(row_count, col_count, x, y, w, h).table
        if column_widths:
            for index, ratio in enumerate(column_widths[:col_count]):
                table.columns[index].width = int(w * float(ratio))
        if row_heights:
            for index, ratio in enumerate(row_heights[:row_count]):
                table.rows[index].height = int(h * float(ratio))
        occupied: set[Tuple[int, int]] = set()
        style = self.table_style(element)
        for row_index, source_row in enumerate(rows_data[:row_count]):
            col_index = 0
            for raw_cell in source_row if isinstance(source_row, list) else []:
                while (row_index, col_index) in occupied and col_index < col_count:
                    col_index += 1
                if col_index >= col_count:
                    break
                cell_data = _as_dict(raw_cell)
                row_span = max(1, int(cell_data.get("rowSpan", 1) or 1))
                col_span = max(1, int(cell_data.get("colSpan", 1) or 1))
                end_row = min(row_count - 1, row_index + row_span - 1)
                end_col = min(col_count - 1, col_index + col_span - 1)
                cell = table.cell(row_index, col_index)
                if end_row != row_index or end_col != col_index:
                    cell = cell.merge(table.cell(end_row, end_col))
                    for rr in range(row_index, end_row + 1):
                        for cc in range(col_index, end_col + 1):
                            if (rr, cc) != (row_index, col_index):
                                occupied.add((rr, cc))
                merged = self.cell_style(style, cell_data, row_index, row_count)
                frame = cell.text_frame
                frame.clear()
                paragraph = frame.paragraphs[0]
                paragraph.alignment = {
                    "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
                }.get(str((merged.get("align") or ["left"])[0]), PP_ALIGN.LEFT)
                run = paragraph.add_run()
                run.text = re.sub(r"<[^>]+>", "", str(cell_data.get("text", "")))
                text_style = self.style({**merged, "style": merged.get("textStyle")})
                self.apply_font(run.font, text_style)
                cell.vertical_anchor = {
                    "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM,
                }.get(str((merged.get("align") or ["left", "middle"])[1]), MSO_ANCHOR.MIDDLE)
                self.apply_fill(cell, merged.get("fill") or element.get("fill"), str(element.get("elementId")))
                col_index += col_span

    def add_chart(self, slide: Any, element: Dict[str, Any]) -> None:
        data = _as_dict(element.get("data"))
        cols = data.get("cols") if isinstance(data.get("cols"), list) else []
        rows = data.get("rows") if isinstance(data.get("rows"), list) else []
        series = element.get("series") if isinstance(element.get("series"), list) else []
        if not cols or not rows or not series:
            self.add_missing_element(slide, element, "图表数据为空")
            return
        first_type = str(_as_dict(series[0]).get("type", "bar"))
        chart_types = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "line": XL_CHART_TYPE.LINE_MARKERS,
                       "area": XL_CHART_TYPE.AREA, "pie": XL_CHART_TYPE.PIE}
        if first_type not in chart_types or any(str(_as_dict(item).get("type")) != first_type for item in series):
            self.add_chart_fallback(slide, element, first_type)
            return
        chart_data = CategoryChartData()
        enc0 = _as_dict(_as_dict(series[0]).get("encode"))
        category_name = enc0.get("item") or enc0.get("x") or cols[0]
        try:
            category_index = cols.index(category_name)
        except ValueError:
            category_index = 0
        chart_data.categories = [str(row[category_index]) for row in rows]
        for index, config_raw in enumerate(series):
            config = _as_dict(config_raw)
            encode = _as_dict(config.get("encode"))
            value_name = encode.get("value") or encode.get("y")
            if value_name not in cols:
                value_name = cols[min(index + 1, len(cols) - 1)]
            value_index = cols.index(value_name)
            values = []
            for row in rows:
                value = row[value_index]
                try:
                    values.append(float(value) if value is not None else None)
                except (TypeError, ValueError):
                    values.append(None)
            chart_data.add_series(str(config.get("name") or value_name), values)
        x, y, w, h = self.rect(element["bounds"])
        chart = slide.shapes.add_chart(chart_types[first_type], x, y, w, h, chart_data).chart
        chart.has_legend = bool(element.get("legend", len(series) > 1))
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        title = element.get("title")
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = str(_as_dict(title).get("text") if isinstance(title, dict) else title)
        for index, output_series in enumerate(chart.series):
            config = _as_dict(series[index])
            color = config.get("fill") or config.get("lineColor")
            if isinstance(color, list):
                color = color[0] if color else None
            if isinstance(color, dict):
                color = _as_dict((color.get("stops") or [{}])[0]).get("color")
            if color:
                output_series.format.fill.solid()
                output_series.format.fill.fore_color.rgb = _hex_color(color, self.colors)
                output_series.format.line.color.rgb = _hex_color(color, self.colors)

    def add_chart_fallback(self, slide: Any, element: Dict[str, Any], chart_type: str) -> None:
        self.warnings.append(f"{element.get('elementId')}: chart type {chart_type!r} rendered as a data summary")
        data = _as_dict(element.get("data"))
        cols = data.get("cols") if isinstance(data.get("cols"), list) else []
        rows = data.get("rows") if isinstance(data.get("rows"), list) else []
        preview = [" · ".join(str(value) for value in cols)]
        preview.extend(" · ".join("—" if value is None else str(value) for value in row) for row in rows[:8])
        surrogate = dict(element)
        surrogate["content"] = {
            "fontSize": 13, "color": "$text", "text": "\n".join(preview), "align": ["left", "top"]
        }
        self.add_text(slide, surrogate)

    def add_icon(self, slide: Any, element: Dict[str, Any]) -> None:
        x, y, w, h = self.rect(element["bounds"])
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        self.apply_fill(shape, element.get("fill") or {"type": "solid", "color": "$primary"}, str(element.get("elementId")))
        shape.line.fill.background()
        frame = shape.text_frame
        frame.clear()
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        icon_name = str(element.get("iconName", "icon")).split(":")[-1]
        run.text = icon_name[:1].upper()
        self.apply_font(run.font, {"fontFamily": "Arial", "fontSize": max(10, min(28, h / 12700)), "bold": True, "color": "#FFFFFF"})
        self.warnings.append(f"{element.get('elementId')}: Font Awesome icon represented by a labelled local glyph")

    def add_missing_element(self, slide: Any, element: Dict[str, Any], label: str) -> None:
        surrogate = dict(element)
        surrogate["content"] = {"fontSize": 12, "color": "#9F1239", "text": f"[{label}]", "align": ["center", "middle"]}
        self.add_text(slide, surrogate)
        self.warnings.append(f"{element.get('elementId')}: {label}")

    def warn_unsupported_features(self, element: Dict[str, Any]) -> None:
        label = str(element.get("elementId") or "element")
        if element.get("opacity") not in (None, 1, 1.0):
            self.warnings.append(f"{label}: element opacity is not supported by the local renderer")
        if any(element.get("flip") or []):
            self.warnings.append(f"{label}: element flip is not supported by the local renderer")
        if element.get("shadow"):
            self.warnings.append(f"{label}: shadow is not supported by the local renderer")
        element_type = element.get("elementType")
        if element_type == "text":
            content = _as_dict(element.get("content"))
            for feature in ("gradient", "shadow", "backgroundColor", "letterSpacing"):
                if content.get(feature) is not None:
                    self.warnings.append(f"{label}: text {feature} is not supported locally")
        elif element_type == "shape" and element.get("adjustments"):
            self.warnings.append(f"{label}: custom shape adjustments use the PowerPoint default locally")
        elif element_type == "image":
            for feature in ("cropShape", "border", "shadow"):
                if element.get(feature) is not None:
                    self.warnings.append(f"{label}: image {feature} is not supported locally")
        elif element_type == "table":
            rows = element.get("rows") if isinstance(element.get("rows"), list) else []
            if any(_as_dict(cell).get("border") is not None for row in rows if isinstance(row, list) for cell in row):
                self.warnings.append(f"{label}: per-cell table borders use PowerPoint defaults locally")

    def apply_background(self, slide: Any, background: Any, page_name: str) -> None:
        spec = _as_dict(background) or {"type": "solid", "color": "#FFFFFF"}
        if spec.get("type", "solid") == "solid":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _hex_color(spec.get("color"), self.colors, "#FFFFFF")
        else:
            slide.background.fill.solid()
            stops = spec.get("stops") if isinstance(spec.get("stops"), list) else []
            color = _as_dict(stops[0]).get("color") if stops else "#FFFFFF"
            slide.background.fill.fore_color.rgb = _hex_color(color, self.colors, "#FFFFFF")
            self.warnings.append(f"{page_name}: complex background approximated with a solid color")

    def export(self, output: Path, transition: str, force: bool) -> Dict[str, Any]:
        output = output.expanduser().resolve()
        output.parent.mkdir(parents=True, exist_ok=True)
        if output.exists() and not force:
            raise ExportError(f"output already exists (pass --force to replace it): {output}")
        prs = Presentation()
        prs.slide_width = Inches(self.slide_width_in)
        prs.slide_height = Inches(self.slide_height_in)
        prs.core_properties.title = str(self.deck.get("title") or self.manifest.stem)
        blank = prs.slide_layouts[6]
        handlers = {
            "text": self.add_text, "shape": self.add_shape, "line": self.add_line,
            "image": self.add_image, "table": self.add_table, "chart": self.add_chart,
            "icon": self.add_icon,
        }
        try:
            for entry in self.deck.get("pages", []):
                page_path = (self.root / str(entry)).resolve()
                try:
                    page_path.relative_to(self.root)
                except ValueError:
                    raise ExportError(f"page path escapes the deck root: {entry}")
                _page_text, page = read_yaml_mapping(page_path)
                slide = prs.slides.add_slide(blank)
                self.apply_background(slide, page.get("background"), str(entry))
                if page.get("notes"):
                    self.warnings.append(f"{entry}: speaker notes are not emitted by the local renderer")
                if page.get("animations"):
                    self.warnings.append(f"{entry}: element animations are not emitted by the local renderer")
                for element in page.get("elements", []):
                    record = _as_dict(element)
                    self.warn_unsupported_features(record)
                    handler = handlers.get(str(record.get("elementType")))
                    if handler is None:
                        self.add_missing_element(slide, record, f"不支持的元素 {record.get('elementType')}")
                    else:
                        try:
                            handler(slide, record)
                        except Exception as exc:
                            self.add_missing_element(slide, record, "本地渲染失败")
                            self.warnings.append(f"{record.get('elementId')}: {type(exc).__name__}: {exc}")
            prs.save(str(output))
        finally:
            for path in self._temp_files:
                path.unlink(missing_ok=True)
        patched = patch_transitions(output, transition)
        result = verify_output(output, transition, expect_fonts=False)
        result.update({
            "transitionPatchedSlides": patched,
            "output": str(output),
            "exporter": "local-python-pptx",
            "warnings": list(dict.fromkeys(self.warnings)),
        })
        return result


def export_local_pptx(manifest: Path, output: Path, transition: str, force: bool = False) -> Dict[str, Any]:
    if PPTX_IMPORT_ERROR is not None:  # pragma: no cover - depends on host environment
        raise ExportError(
            "local PPTX export requires python-pptx; install it with "
            "'python3 -m pip install python-pptx'"
        ) from PPTX_IMPORT_ERROR
    return LocalPptdExporter(manifest).export(output, transition, force)
