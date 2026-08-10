#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import shutil
import tempfile
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

UNSUPPORTED_MARKER = "LEGALWORK_DOCUMENT_UNSUPPORTED"

try:
    from pptx import Presentation
except Exception as exc:  # pragma: no cover - environment-dependent
    Presentation = None  # type: ignore[assignment]
    IMPORT_ERROR = str(exc)
else:
    IMPORT_ERROR = ""


def emit(payload: dict[str, Any], exit_code: int = 0) -> None:
    print(json.dumps(payload, ensure_ascii=False, separators=(",", ":")))
    raise SystemExit(exit_code)


def unsupported(reason: str, operation: str, detail: Any | None = None) -> None:
    root = Path(tempfile.gettempdir()) / "legalwork-office-fallback"
    root.mkdir(parents=True, exist_ok=True)
    ticket = root / f"ticket-{uuid.uuid4().hex}.json"
    ticket.write_text(json.dumps({
        "marker": UNSUPPORTED_MARKER,
        "status": "unsupported",
        "source": "legal-document-formatting",
        "operation": operation,
        "reason": reason,
        "detail": detail,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }, ensure_ascii=False), encoding="utf-8")
    payload = {
        "status": "unsupported",
        "marker": UNSUPPORTED_MARKER,
        "operation": operation,
        "reason": reason,
        "fallback": "office_mcp",
        "fallback_ticket": str(ticket),
    }
    if detail is not None:
        payload["detail"] = detail
    emit(payload)


def fail(message: str, operation: str) -> None:
    emit({"status": "error", "operation": operation, "error": message}, 1)


def require_pptx(operation: str) -> None:
    if Presentation is None:
        unsupported(f"python-pptx unavailable: {IMPORT_ERROR}", operation)


def ensure_input(path: str, operation: str) -> Path:
    p = Path(path).expanduser().resolve()
    if not p.is_file():
        fail(f"file not found: {p}", operation)
    if p.suffix.lower() != ".pptx":
        unsupported(f"only .pptx is handled locally (got {p.suffix or 'no extension'})", operation)
    return p


def copy_output(src: Path, output: str, operation: str) -> Path:
    dst = Path(output).expanduser().resolve()
    if dst.suffix.lower() != ".pptx":
        fail("output must end with .pptx", operation)
    dst.parent.mkdir(parents=True, exist_ok=True)
    if src != dst:
        shutil.copy2(src, dst)
    return dst


def slide_text(slide: Any) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def cmd_inspect(args: argparse.Namespace) -> None:
    operation = "inspect"
    require_pptx(operation)
    path = ensure_input(args.input, operation)
    try:
        prs = Presentation(str(path))
        samples = []
        for index, slide in enumerate(prs.slides, start=1):
            if index > 8:
                break
            text = slide_text(slide)
            samples.append({"slide": index, "text": text[:240]})
        emit({
            "status": "ok",
            "operation": operation,
            "summary": {
                "file": str(path),
                "slides": len(prs.slides),
                "size_inches": [round(prs.slide_width / 914400, 2), round(prs.slide_height / 914400, 2)],
                "slide_sample": samples,
            },
        })
    except Exception as exc:
        fail(str(exc), operation)


def cmd_replace(args: argparse.Namespace) -> None:
    operation = "replace"
    require_pptx(operation)
    src = ensure_input(args.input, operation)
    dst = copy_output(src, args.output, operation)
    try:
        prs = Presentation(str(dst))
        replaced = 0
        spanning = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if args.old not in paragraph.text:
                        continue
                    local = 0
                    for run in paragraph.runs:
                        count = run.text.count(args.old)
                        if count:
                            run.text = run.text.replace(args.old, args.new)
                            local += count
                    if local:
                        replaced += local
                    else:
                        spanning += 1
        if spanning:
            unsupported(
                "target text crosses PowerPoint run boundaries; safe local replacement would not preserve formatting",
                operation,
                {"spanning_paragraphs": spanning, "single_run_replacements": replaced},
            )
        if replaced:
            prs.save(str(dst))
        emit({"status": "ok", "operation": operation, "output": str(dst), "replacements": replaced})
    except SystemExit:
        raise
    except Exception as exc:
        fail(str(exc), operation)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="LegalWork deterministic PPTX worker")
    sub = parser.add_subparsers(dest="command", required=True)
    p = sub.add_parser("inspect")
    p.add_argument("--input", required=True)
    p.set_defaults(func=cmd_inspect)
    p = sub.add_parser("replace")
    p.add_argument("--input", required=True)
    p.add_argument("--output", required=True)
    p.add_argument("--old", required=True)
    p.add_argument("--new", required=True)
    p.set_defaults(func=cmd_replace)
    return parser


def main() -> None:
    args = build_parser().parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
